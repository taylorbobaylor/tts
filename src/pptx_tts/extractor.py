"""Content extraction from PowerPoint files."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


@dataclass
class SlideContent:
    """Text content extracted from a single slide."""

    number: int
    title: str
    body: list[str] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        """Return all text for this slide as a single readable string."""
        parts: list[str] = []
        if self.title:
            parts.append(self.title)
        parts.extend(line for line in self.body if line.strip())
        return ". ".join(parts)


def extract_slides(filepath: str | Path) -> list[SlideContent]:
    """Extract text content from every slide in a .pptx file.

    Args:
        filepath: Path to a .pptx file.

    Returns:
        Ordered list of SlideContent, one per slide.

    Raises:
        FileNotFoundError: If the file does not exist.
        ValueError: If the file is not a .pptx file.
    """
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Expected a .pptx file, got: {path.suffix}")

    prs = Presentation(str(path))
    slides: list[SlideContent] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        body_lines: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Use the first title-shaped placeholder as the title
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title = text
            else:
                body_lines.append(text)

        slides.append(SlideContent(number=idx, title=title, body=body_lines))

    return slides
