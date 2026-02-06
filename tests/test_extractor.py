"""Tests for the content extractor module."""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_tts.extractor import SlideContent, extract_slides


def _make_pptx(slides_data: list[tuple[str, str]], path: Path) -> Path:
    """Helper: create a .pptx with the given (title, body) per slide."""
    prs = Presentation()
    for title, body in slides_data:
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)
        slide.placeholders[0].text = title
        slide.placeholders[1].text = body
    prs.save(str(path))
    return path


class TestExtractSlides:
    def test_single_slide(self, tmp_path: Path) -> None:
        pptx = _make_pptx([("Hello", "World")], tmp_path / "one.pptx")
        slides = extract_slides(pptx)
        assert len(slides) == 1
        assert slides[0].title == "Hello"
        assert slides[0].body == ["World"]

    def test_multiple_slides(self, tmp_path: Path) -> None:
        data = [("Slide 1", "Body 1"), ("Slide 2", "Body 2"), ("Slide 3", "Body 3")]
        pptx = _make_pptx(data, tmp_path / "multi.pptx")
        slides = extract_slides(pptx)
        assert len(slides) == 3
        assert [s.title for s in slides] == ["Slide 1", "Slide 2", "Slide 3"]

    def test_empty_presentation(self, tmp_path: Path) -> None:
        pptx = tmp_path / "empty.pptx"
        Presentation().save(str(pptx))
        slides = extract_slides(pptx)
        assert slides == []

    def test_file_not_found(self) -> None:
        with pytest.raises(FileNotFoundError):
            extract_slides("/nonexistent/file.pptx")

    def test_wrong_extension(self, tmp_path: Path) -> None:
        bad = tmp_path / "file.txt"
        bad.write_text("not a pptx")
        with pytest.raises(ValueError, match=".pptx"):
            extract_slides(bad)


class TestSlideContent:
    def test_full_text_title_and_body(self) -> None:
        sc = SlideContent(number=1, title="Title", body=["Line 1", "Line 2"])
        assert sc.full_text == "Title. Line 1. Line 2"

    def test_full_text_no_title(self) -> None:
        sc = SlideContent(number=1, title="", body=["Only body"])
        assert sc.full_text == "Only body"

    def test_full_text_empty(self) -> None:
        sc = SlideContent(number=1, title="", body=[])
        assert sc.full_text == ""
