"""Tests for the playback controller module."""

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation

from pptx_tts.extractor import SlideContent
from pptx_tts.playback import ENDING_JOKES, PlaybackController


def _make_pptx(slides_data: list[tuple[str, str]], path: Path) -> Path:
    prs = Presentation()
    for title, body in slides_data:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.placeholders[0].text = title
        slide.placeholders[1].text = body
    prs.save(str(path))
    return path


@patch("pptx_tts.playback.Synthesizer")
class TestPlaybackController:
    def test_reads_all_slides(self, mock_synth_cls: MagicMock, tmp_path: Path) -> None:
        pptx = _make_pptx([("S1", "B1"), ("S2", "B2")], tmp_path / "test.pptx")
        mock_synth = mock_synth_cls.return_value

        ctrl = PlaybackController(slide_delay=0, tell_joke=False)
        ctrl.play_presentation(str(pptx))

        # Two slides spoken
        assert mock_synth.speak.call_count == 2
        spoken_texts = [call.args[0] for call in mock_synth.speak.call_args_list]
        assert "S1" in spoken_texts[0]
        assert "S2" in spoken_texts[1]

    def test_ending_joke_spoken(self, mock_synth_cls: MagicMock, tmp_path: Path) -> None:
        pptx = _make_pptx([("Slide", "Content")], tmp_path / "joke.pptx")
        mock_synth = mock_synth_cls.return_value

        ctrl = PlaybackController(slide_delay=0, tell_joke=True)
        ctrl.play_presentation(str(pptx))

        # Last speak call should be one of the jokes
        last_text = mock_synth.speak.call_args_list[-1].args[0]
        assert last_text in ENDING_JOKES

    def test_stop_halts_playback(self, mock_synth_cls: MagicMock, tmp_path: Path) -> None:
        pptx = _make_pptx(
            [("S1", "B1"), ("S2", "B2"), ("S3", "B3")], tmp_path / "stop.pptx"
        )
        mock_synth = mock_synth_cls.return_value

        ctrl = PlaybackController(slide_delay=0, tell_joke=False)
        # Stop immediately after first speak
        mock_synth.speak.side_effect = lambda _: ctrl.stop()
        ctrl.play_presentation(str(pptx))

        # Only one slide should have been spoken
        assert mock_synth.speak.call_count == 1
